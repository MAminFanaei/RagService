# test_quality_pptx.py

import json
import time
from pathlib import Path

def save_result(output_path, method, data):
    with open(output_path / f"{method}.json", "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"  ✅ {method:<20} → {output_path / f'{method}.json'}")

def run_pptx_tests(file_path: Path, output_dir: Path, args):
    print(f"📊 Testing PPTX parser on: {file_path.name}\n")
    
    print("▶ Testing python-pptx...")
    try:
        from pptx import Presentation
        
        t0 = time.perf_counter()
        prs = Presentation(file_path)
        
        slides_data = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            text_parts = []
            notes = ""
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
            
            # Extract speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    notes = notes_frame.text.strip()
            
            slides_data.append({
                "slide_number": slide_num,
                "text": "\n".join(text_parts),
                "notes": notes,
                "shape_count": len(slide.shapes),
            })
        
        elapsed = time.perf_counter() - t0
        
        result = {
            "parser": "python-pptx",
            "file": file_path.name,
            "time_seconds": round(elapsed, 3),
            "slide_count": len(slides_data),
            "slides_preview": slides_data[:args.max_chunks],
            "verdict": {
                "quality": f"✅ {len(slides_data)} slides extracted",
            }
        }
        save_result(output_dir, "parser_python_pptx", result)
        
    except ImportError:
        print(f"  ⏭️  python-pptx not installed — skip")
    except Exception as e:
        print(f"  ❌ python-pptx failed: {e}")