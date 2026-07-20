"""
PPTX parser using python-pptx.

Extracts:
  - Slide title → heading element
  - Text frames → text elements
  - Tables → Markdown via tabulate
  - Speaker notes → text elements
  - Slide number = page number
"""

from __future__ import annotations

import io
from typing import Any

import structlog

from ingestion.parsers.base import BaseParser, ParsedElement, ParserError

logger = structlog.get_logger(__name__)


def _shape_text(shape) -> str:
    """Safely extract text from a pptx shape."""
    try:
        if shape.has_text_frame:
            return "\n".join(
                para.text.strip()
                for para in shape.text_frame.paragraphs
                if para.text.strip()
            )
    except Exception:
        pass
    return ""


def _table_to_markdown(table) -> tuple[str, str]:
    """Return (plain_text, markdown) for a pptx Table."""
    try:
        from tabulate import tabulate  # type: ignore

        rows = []
        for row in table.rows:
            rows.append([cell.text_frame.text.strip() for cell in row.cells])

        if not rows:
            return "", ""

        plain = "\n".join(" | ".join(r) for r in rows)
        if len(rows) > 1:
            md = tabulate(rows[1:], headers=rows[0], tablefmt="pipe")
        else:
            md = tabulate(rows, tablefmt="pipe")
        return plain, md

    except ImportError:
        rows = []
        try:
            for row in table.rows:
                rows.append([cell.text_frame.text.strip() for cell in row.cells])
        except Exception:
            return "", ""
        plain = "\n".join(" | ".join(r) for r in rows)
        if len(rows) > 1:
            header = " | ".join(rows[0])
            sep = " | ".join(["---"] * len(rows[0]))
            body = "\n".join(" | ".join(r) for r in rows[1:])
            md = f"{header}\n{sep}\n{body}"
        else:
            md = plain
        return plain, md


class PptxParser(BaseParser):
    """Parser for PowerPoint .pptx files."""

    @property
    def name(self) -> str:
        return "pptx"

    async def parse(self, data: bytes, filename: str) -> list[ParsedElement]:
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Emu  # type: ignore
        except ImportError as exc:
            raise ParserError(self.name, "python-pptx not installed", exc)

        try:
            prs = Presentation(io.BytesIO(data))
        except Exception as exc:
            raise ParserError(self.name, f"Cannot open PPTX: {exc}", exc)

        elements: list[ParsedElement] = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            slide_title: str | None = None

            # ---- Title shape ----
            try:
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    slide_title = slide.shapes.title.text_frame.text.strip() or None
            except Exception:
                pass

            if slide_title:
                elements.append(
                    ParsedElement(
                        text=slide_title,
                        parser_name=self.name,
                        element_type="heading",
                        heading_level=1,
                        page_number=slide_num,
                        section_title=slide_title,
                        section_path=[slide_title],
                        raw_metadata={"slide": slide_num},
                    )
                )

            # ---- Other shapes ----
            for shape in slide.shapes:
                # Skip the title shape (already handled)
                try:
                    if slide.shapes.title and shape == slide.shapes.title:
                        continue
                except Exception:
                    pass

                shape_type = getattr(shape, "shape_type", None)

                # Table
                if shape.has_table if hasattr(shape, "has_table") else False:
                    try:
                        plain, md = _table_to_markdown(shape.table)
                        if plain.strip():
                            elements.append(
                                ParsedElement(
                                    text=plain,
                                    parser_name=self.name,
                                    element_type="table",
                                    is_table=True,
                                    table_markdown=md,
                                    page_number=slide_num,
                                    section_title=slide_title,
                                    section_path=[slide_title] if slide_title else [],
                                    raw_metadata={"slide": slide_num},
                                )
                            )
                    except Exception:
                        pass
                    continue

                # Text frame
                text = _shape_text(shape)
                if text:
                    elements.append(
                        ParsedElement(
                            text=text,
                            parser_name=self.name,
                            element_type="text",
                            page_number=slide_num,
                            section_title=slide_title,
                            section_path=[slide_title] if slide_title else [],
                            raw_metadata={"slide": slide_num},
                        )
                    )

            # ---- Speaker notes ----
            try:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        elements.append(
                            ParsedElement(
                                text=notes_text,
                                parser_name=self.name,
                                element_type="text",
                                page_number=slide_num,
                                section_title=slide_title,
                                section_path=[slide_title] if slide_title else [],
                                raw_metadata={"slide": slide_num, "is_notes": True},
                            )
                        )
            except Exception:
                pass

        logger.info(
            "pptx_parsed",
            filename=filename,
            slide_count=len(prs.slides),
            element_count=len(elements),
        )
        return elements
